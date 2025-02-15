from pptx import Presentation 
from pptx.util import Inches
import pandas as pd
import barcode
from barcode.writer import ImageWriter
import os

# add the the serial number you need in the first column 
excel_file = "serial_numbers.xlsx"

df = pd.read_excel(excel_file)

prs = Presentation()

# add image to be a background in the current folder
image_path = "gift_card_barcode.jpeg"


for index, row in df.iterrows():
    
    barcode_value = str(row[0])
    barcode_class = barcode.get_barcode_class("code128")
    barcode_image = barcode_class(barcode_value, writer=ImageWriter())

    barcode_filename = f"barcode_{index}"
    barcode_image.save(barcode_filename)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)

    left = Inches(0.5)
    top = slide_height - Inches(1.5)



    slide.shapes.add_picture(f"{barcode_filename}.png", left, top, width=Inches(2), height=Inches(1))

    os.remove(f"{barcode_filename}.png")

    prs.save('card_side_02.pptx')

    print(f"{index}")

    


